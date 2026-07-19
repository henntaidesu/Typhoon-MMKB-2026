# -*- coding: utf-8 -*-
"""Generate the Japanese system-presentation deck for Typhoon MMKB."""
import io, sys
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_LINE_DASH_STYLE

OUT = r"d:/Project/Musashino/ADVANCED-MMKB-2026/docs/MMKB_presentation_ja.pptx"

NAVY   = RGBColor(0x14, 0x2B, 0x4A)
BLUE   = RGBColor(0x1F, 0x5C, 0x99)
ACCENT = RGBColor(0x00, 0x8C, 0x8C)
ORANGE = RGBColor(0xD1, 0x66, 0x1E)
GREY   = RGBColor(0x5A, 0x66, 0x77)
LIGHT  = RGBColor(0xF4, 0xF6, 0xFA)
BORDER = RGBColor(0xC9, 0xD3, 0xE4)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
CODEBG = RGBColor(0xF7, 0xF8, 0xFB)
CODEFG = RGBColor(0x18, 0x2A, 0x3E)

JP = "Meiryo"
MONO = "Consolas"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = 13.333, 7.5
BLANK = prs.slide_layouts[6]

_n = [0]
WARN = []


def rect(slide, x, y, w, h, fill=None, line=None, shape=MSO_SHAPE.RECTANGLE, dash=None):
    s = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line; s.line.width = Pt(1)
        if dash:
            s.line.dash_style = dash
    s.shadow.inherit = False
    s.text_frame.text = ""
    return s


def tb(slide, x, y, w, h, wrap=True):
    t = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = t.text_frame
    tf.word_wrap = wrap
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    return tf


def para(tf, text, size=14, bold=False, color=NAVY, font=JP, space_after=6,
         first=False, align=PP_ALIGN.LEFT, space_before=0, line=None):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    p.space_before = Pt(space_before)
    if line:
        p.line_spacing = line
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = color; r.font.name = font
    return p


def slide(title, sub=None, band=NAVY):
    s = prs.slides.add_slide(BLANK)
    _n[0] += 1
    rect(s, 0, 0, SW, 0.92, fill=band)
    rect(s, 0, 0.92, SW, 0.055, fill=ACCENT)
    tf = tb(s, 0.55, 0.17, 11.6, 0.6)
    para(tf, title, size=25, bold=True, color=WHITE, first=True, space_after=0)
    if sub:
        tf2 = tb(s, 0.58, 0.62, 11.6, 0.3)
        para(tf2, sub, size=11.5, color=RGBColor(0xA8, 0xC4, 0xE0), first=True, space_after=0)
    # slide number
    nf = tb(s, 12.4, 6.95, 0.6, 0.3)
    para(nf, str(_n[0]), size=11, color=GREY, first=True, align=PP_ALIGN.RIGHT, space_after=0)
    return s


def code(s, x, y, w, h, lines, size=11, title=None):
    """Code panel with optional caption strip."""
    if title:
        cap = rect(s, x, y, w, 0.30, fill=NAVY)
        ctf = cap.text_frame
        ctf.margin_left = Inches(0.12); ctf.margin_top = Emu(0); ctf.margin_bottom = Emu(0)
        ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
        para(ctf, title, size=10.5, bold=True, color=WHITE, first=True, space_after=0)
        y += 0.30; h -= 0.30
    rect(s, x, y, w, h, fill=CODEBG, line=BORDER)
    # auto-fit: shrink the font if the lines would overflow the panel
    avail = h - 0.22
    per = size * 1.2 * 1.06 / 72.0
    if lines and len(lines) * per > avail:
        size = round(avail / len(lines) * 72.0 / (1.2 * 1.06), 1)
        WARN.append("shrunk to %.1fpt (%d lines, %.2fin): %r" % (size, len(lines), h, (title or lines[0])[:40]))
    # also guard horizontal overflow (Consolas advance ~0.60 em)
    maxw = max((len(l) for l in lines), default=0) * size * 0.60 / 72.0
    if maxw > w - 0.28:
        size = round((w - 0.28) * 72.0 / (max(len(l) for l in lines) * 0.60), 1)
        WARN.append("narrowed to %.1fpt (%d cols, %.2fin): %r" % (size, max(len(l) for l in lines), w, (title or lines[0])[:40]))
    tf = tb(s, x + 0.14, y + 0.10, w - 0.28, h - 0.2, wrap=False)
    for i, ln in enumerate(lines):
        c = CODEFG
        st = ln.lstrip()
        if st.startswith("--"):
            c = RGBColor(0x6E, 0x8A, 0x5E)          # comment green
        elif st.startswith("=>") or st.startswith("→"):
            c = ORANGE
        para(tf, ln if ln else " ", size=size, color=c, font=MONO,
             first=(i == 0), space_after=0, line=1.06)


def bullets(s, x, y, w, h, items, size=14.5, gap=9):
    tf = tb(s, x, y, w, h)
    first = True
    for it in items:
        if isinstance(it, tuple):
            txt, lvl = it
        else:
            txt, lvl = it, 0
        if lvl == 0:
            para(tf, "▪ " + txt, size=size, bold=True, color=NAVY,
                 first=first, space_after=3, space_before=(0 if first else gap))
        else:
            para(tf, "    " + txt, size=size - 2, color=GREY, first=first, space_after=3)
        first = False
    return tf


def shot(s, x, y, w, h, label="Navicat スクリーンショット挿入位置"):
    r = rect(s, x, y, w, h, fill=RGBColor(0xFA, 0xFB, 0xFD), line=RGBColor(0x9F, 0xB2, 0xCC),
             dash=MSO_LINE_DASH_STYLE.DASH)
    tf = r.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    para(tf, "📷  " + label, size=12.5, bold=True, color=RGBColor(0x7F, 0x92, 0xAD),
         first=True, align=PP_ALIGN.CENTER, space_after=0)
    return r


def kpi(s, x, y, w, h, value, label, color=BLUE):
    rect(s, x, y, w, h, fill=WHITE, line=BORDER)
    rect(s, x, y, 0.075, h, fill=color)
    tf = tb(s, x + 0.24, y + 0.16, w - 0.35, h - 0.25)
    para(tf, value, size=22, bold=True, color=color, first=True, space_after=1)
    para(tf, label, size=10.5, color=GREY, space_after=0)


def note(s, x, y, w, h, text, color=ACCENT):
    rect(s, x, y, w, h, fill=RGBColor(0xEC, 0xF6, 0xF6), line=color)
    tf = tb(s, x + 0.18, y + 0.12, w - 0.36, h - 0.24)
    para(tf, text, size=12, color=RGBColor(0x0B, 0x4C, 0x4C), first=True, space_after=0, line=1.15)


# ===========================================================================
# 1. Title
# ===========================================================================
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SW, SH, fill=NAVY)
rect(s, 0, 4.62, SW, 0.06, fill=ACCENT)
tf = tb(s, 1.0, 1.65, 11.3, 1.2)
para(tf, "台風マルチメディア知識ベース", size=40, bold=True, color=WHITE, first=True, space_after=6)
para(tf, "Typhoon MMKB — 時空間計算と意味計算を統合した分散知識ベース", size=17,
     color=RGBColor(0x8F, 0xB8, 0xDC), space_after=0)
tf = tb(s, 1.0, 4.95, 11.3, 1.6)
para(tf, "マルチメディア知識ベース構成論   Subject-1 ／ 選題 (8) Disasters",
     size=14, color=RGBColor(0xC9, 0xDA, 0xEA), first=True, space_after=8)
para(tf, "武蔵野大学  Spring 2026", size=13, color=RGBColor(0x7F, 0x9C, 0xBB), space_after=8)
para(tf, "PostgreSQL 18.4 + PostGIS 3.6.2 + pgvector 0.8.5  ／  FastAPI  ／  Vue 3 + Leaflet",
     size=12, color=RGBColor(0x6A, 0x8A, 0xAA), space_after=0)

# ===========================================================================
# 2. 設計思路 — 課題と三層計算モデル
# ===========================================================================
s = slide("1. システム設計思路 — 三層計算モデル", "台風情報は分散・多言語・非構造。語句一致でも属性検索でも答えられない問いに応える")

rect(s, 0.75, 1.3, 11.8, 0.72, fill=LIGHT, line=BORDER)
tf = tb(s, 1.0, 1.42, 11.3, 0.5)
para(tf, "課題   経路は NOAA、被害は GDACS、警報は気象庁・中央气象台… 形式も言語 (JP/EN/CN) もバラバラ",
     size=12.5, bold=True, color=NAVY, first=True, space_after=3)
para(tf, "解    分散ソースを 1 つの知識ベースへ統合し、空間・時間・意味の 3 軸で「選択」し「結合」する",
     size=12.5, bold=True, color=BLUE, space_after=0)

cols = [
    ("空間的選択", "Spatial Selection", "ST_Intersects(geom, bbox)",
     "PostGIS  GiST 索引", "track_point.geom", BLUE),
    ("時間的選択", "Temporal Selection", "obs_time BETWEEN ? AND ?",
     "B-tree 索引", "obs_time / event_time", ACCENT),
    ("意味的選択", "Semantic Selection", "embedding <=> query_vec",
     "pgvector  IVFFlat 索引", "384 次元 / JP·EN·CN 共有", ORANGE),
]
x = 0.75
for t, e, sql, idx, col, c in cols:
    rect(s, x, 2.2, 3.85, 2.15, fill=WHITE, line=BORDER)
    rect(s, x, 2.2, 3.85, 0.09, fill=c)
    tf = tb(s, x + 0.28, 2.44, 3.3, 1.8)
    para(tf, t, size=16, bold=True, color=c, first=True, space_after=1)
    para(tf, e, size=9.5, color=GREY, space_after=11)
    para(tf, sql, size=10, color=CODEFG, font=MONO, space_after=8)
    para(tf, "索引: " + idx, size=10, color=GREY, space_after=3)
    para(tf, "対象: " + col, size=10, color=GREY, space_after=0)
    x += 4.05

rect(s, 2.6, 4.55, 8.15, 0.85, fill=NAVY)
tf = tb(s, 2.9, 4.72, 7.6, 0.6)
para(tf, "結合 (Join) — 時空間で候補を絞り、意味距離で並べ替える", size=15, bold=True,
     color=WHITE, first=True, align=PP_ALIGN.CENTER, space_after=2)
para(tf, "/search/hybrid  →  「意味的結合」の工学的実装", size=10.5,
     color=RGBColor(0x9F, 0xC4, 0xE2), align=PP_ALIGN.CENTER, space_after=0)

note(s, 0.75, 5.6, 11.8, 1.15,
     "多言語モデル paraphrase-multilingual-MiniLM-L12-v2 により JP / EN / CN が同一の 384 次元意味空間を共有する。\n"
     "検証済み: 「洪水と土砂災害をもたらした台風」と \"severe flooding typhoon\" の類似度 0.82、無関係文とは 0.10。\n"
     "英語クエリ \"severe flooding and heavy rain warning\" が中国語「暴雨预警信号」を距離 0.33 で命中する。")

# ===========================================================================
# 3. システム構成 + 実装規模
# ===========================================================================
s = slide("1. システム設計思路 — 構成と実装規模", "分散ソース → 統合 → 知識ベース → API → UI ／ 実データは 2026-07 現在")
layers = [
    ("分散データソース", "IBTrACS(NOAA) ／ GDACS(UN/EC) ／ Digital Typhoon(NII) ／ Natural Earth ／ GADM 4.1\n"
                        "中央气象台・香港天文台・気象庁警報 ／ 応急管理部・消防庁・ReliefWeb", RGBColor(0x4A, 0x6C, 0x8E)),
    ("クローラ / 統合  backend/crawler/", "load.py 三段階マッチング (intl_id → 名称 → 時空間近傍)\n"
                        "enrich.py 地理影響の導出 ／ embed.py 多言語ベクトル化 ／ pipeline.py 統括", BLUE),
    ("知識ベース  PostgreSQL 18.4", "9 テーブル ／ PostGIS 3.6.2 + pgvector 0.8.5\n"
                        "トリガ 7 本・関数 6 本 ／ GiST + IVFFlat + B-tree 索引", ACCENT),
    ("バックエンド  FastAPI + SQLAlchemy", "/typhoons 属性 ・ /search/spatiotemporal ・ /search/semantic\n"
                        "/search/hybrid 結合 ・ /stats/* 集計 — すべて GeoJSON で応答", RGBColor(0x2E, 0x7D, 0x5B)),
    ("フロントエンド  Vue 3 + Leaflet", "地図 (経路・強度配色・被害点・上陸点) ／ 時間軸再生\n"
                        "意味連想検索 ／ 三段階コロプレス統計", ORANGE),
]
y = 1.32
for t, d, c in layers:
    rect(s, 0.75, y, 7.15, 0.98, fill=WHITE, line=BORDER)
    rect(s, 0.75, y, 0.09, 0.98, fill=c)
    tf = tb(s, 1.02, y + 0.11, 6.75, 0.82)
    para(tf, t, size=12.5, bold=True, color=c, first=True, space_after=3)
    para(tf, d, size=9.5, color=GREY, space_after=0, line=1.1)
    y += 0.98
    if y < 6.2:
        a = tb(s, 4.2, y - 0.05, 0.5, 0.2)
        para(a, "▼", size=8, color=BORDER, first=True, space_after=0)
    y += 0.09

k = [("1,993", "台風 (1949–2026)", BLUE), ("135,228", "経路点", BLUE),
     ("929", "受灾情報", ORANGE), ("210", "公共情報", ORANGE),
     ("5,400", "行政区 (国22/省591/市4787)", ACCENT), ("383,707", "影響事実", ACCENT),
     ("1,905", "上陸イベント", RGBColor(0x2E, 0x7D, 0x5B)), ("158 MB", "DB 総容量", GREY)]
x, y = 8.15, 1.32
for i, (v, l, c) in enumerate(k):
    kpi(s, x, y, 2.1, 0.98, v, l, c)
    x += 2.22
    if i % 2 == 1:
        x = 8.15; y += 1.07

note(s, 8.15, 5.6, 4.4, 1.15,
     "観測機関別  JMA 71,163 ／ CMA 60,498 ／ JTWC 3,567\n"
     "上陸回数 (実測)  China 600 ＞ Philippines 527\n"
     "＞ Japan 284 ＞ Vietnam 250 ＞ Taiwan 106")

# ===========================================================================
# 4. 主なシーン
# ===========================================================================
s = slide("2. 主なシーン", "地図可視化 ／ 意味連想検索 ／ 統計分析 ／ データ取込")
scenes = [
    ("① 地図可視化と時間軸再生", "経路を強度で配色して描画し、時間軸スライダで進行を再生。被害点・公共情報・上陸点を重畳",
     "地図画面 (MapPage.vue)", BLUE),
    ("② 意味連想検索", "自然言語を JP/EN/CN いずれでも入力。台風・受灾情報・公共情報の 3 層を同時に意味検索",
     "意味検索 (SemanticSearchBox.vue)", ORANGE),
    ("③ 統計 — 三段階コロプレス地図", "国 / 省 / 地級市を切替。行政区クリックで影響した台風の経路を地図に描画",
     "統計画面 (StatsView.vue)", ACCENT),
    ("④ データソース管理と取込", "各ソースの取得状況を一覧。Web UI から取込を実行し進捗を監視",
     "データソース画面 (DataSources.vue)", RGBColor(0x2E, 0x7D, 0x5B)),
]
for i, (t, d, cap, c) in enumerate(scenes):
    x = 0.75 + (i % 2) * 6.0
    y = 1.3 + (i // 2) * 2.75
    rect(s, x, y, 5.8, 2.62, fill=WHITE, line=BORDER)
    rect(s, x, y, 5.8, 0.07, fill=c)
    tf = tb(s, x + 0.22, y + 0.2, 5.35, 0.55)
    para(tf, t, size=13.5, bold=True, color=c, first=True, space_after=3)
    para(tf, d, size=9.5, color=GREY, space_after=0, line=1.12)
    shot(s, x + 0.22, y + 0.98, 5.35, 1.5, cap)

# ===========================================================================
# 5. DB 全体像
# ===========================================================================
s = slide("3. データベース設計 — 全体像", "9 テーブル ／ SQLAlchemy ORM (models.py) で定義し create_all で生成")
groups = [
    ("中核 (Core)", [("typhoon", "台風 = 最上位知識単位 + 意味ベクトル", True),
                     ("track_point", "経路点 = 時空間の核 (135,228 行)", True)], BLUE),
    ("情報層 (Information)", [("secondary_disaster", "受灾情報 — 実際に発生した被害 + ベクトル", True),
                              ("public_info", "公共情報 — 当局発表の警報・避難 + ベクトル", False),
                              ("media_asset", "衛星画像・写真のメタデータ", False)], ORANGE),
    ("地理層 (Geography)", [("admin_region", "参照行政境界 国 / 省 / 地級市", False),
                            ("typhoon_region_impact", "台風 × 行政区 の影響事実 (導出)", False),
                            ("landfall", "上陸イベント (導出)", False),
                            ("affected_region", "影響範囲ポリゴン", False)], ACCENT),
]
x = 0.75
for gname, rows, c in groups:
    rect(s, x, 1.32, 3.85, 4.1, fill=WHITE, line=BORDER)
    rect(s, x, 1.32, 3.85, 0.42, fill=c)
    gtf = tb(s, x + 0.22, 1.42, 3.4, 0.3)
    para(gtf, gname, size=13, bold=True, color=WHITE, first=True, space_after=0)
    tf = tb(s, x + 0.22, 1.92, 3.45, 3.35)
    first = True
    for tn, td, star in rows:
        para(tf, ("★ " if star else "・ ") + tn, size=12.5, bold=True,
             color=(ORANGE if star else NAVY), first=first, space_after=2)
        para(tf, "      " + td, size=9.5, color=GREY, space_after=11, line=1.1)
        first = False
    x += 4.05
note(s, 0.75, 5.62, 11.8, 1.1,
     "★ = 本発表で DDL を提示する主要 3 テーブル。台風 (知識単位) → 経路点 (時空間) → 受灾情報 (意味) の順に、\n"
     "設計の 3 本柱をそれぞれ代表している。他テーブルは同じ設計原則の適用例のため割愛する。\n"
     "共通原則: すべての知識単位が「空間列 geom」と「意味列 embedding」を同時に持ち、親台風へ CASCADE で従属する。")

# ===========================================================================
# 6-8. 主要テーブル DDL
# ===========================================================================
s = slide("3. 主要テーブル ① typhoon", "台風 = 最上位の知識単位。構造化属性と意味ベクトルを同居させる")
code(s, 0.75, 1.32, 7.4, 4.15, [
    "CREATE TABLE typhoon (",
    "    id               SERIAL       PRIMARY KEY,",
    "    intl_id          VARCHAR(16)  UNIQUE NOT NULL,  -- 国際番号 例:2306",
    "    sid              VARCHAR(32),                   -- IBTrACS storm id",
    "    name             VARCHAR(64),   name_jp VARCHAR(64),",
    "    name_cn          VARCHAR(64),   season_year INTEGER,",
    "    category         VARCHAR(32),      -- TY / STS / SuperTY",
    "    max_wind_kt      DOUBLE PRECISION, -- トリガ#2 が自動更新",
    "    min_pressure_hpa DOUBLE PRECISION, -- トリガ#2 が自動更新",
    "    start_time       TIMESTAMPTZ,   end_time TIMESTAMPTZ,",
    "    source           VARCHAR(32),   is_active BOOLEAN,",
    "    -- ▼ 意味層: 属性から生成した多言語要約と、その 384 次元ベクトル",
    "    summary_text     TEXT,",
    "    embedding        VECTOR(384)",
    ");",
    "CREATE INDEX ix_typhoon_season_year ON typhoon (season_year);",
    "-- 意味的選択のためのコサイン距離用 IVFFlat 索引",
    "CREATE INDEX ix_typhoon_embedding ON typhoon",
    "    USING ivfflat (embedding vector_cosine_ops) WITH (lists = 100);",
], size=10.5, title="DDL — typhoon  (1,993 行)")
code(s, 0.75, 5.6, 7.4, 1.15, [
    "-- 主要クエリ Q1 : 属性選択  → 2302 Mawar / 2315 Bolaven / 2309 Saola (SuperTY)",
    "SELECT intl_id, name, category, max_wind_kt FROM typhoon",
    " WHERE season_year = 2023 AND max_wind_kt >= 100 ORDER BY max_wind_kt DESC;",
], size=10, title="主要クエリ Q1 — 属性選択 (Attribute Selection)")
tf = tb(s, 8.4, 1.4, 4.2, 5.3)
para(tf, "設計のポイント", size=15, bold=True, color=BLUE, first=True, space_after=12)
for t, d in [("intl_id が冪等キー", "再クロール時の UPSERT はここで一意に解決する"),
             ("summary_text + embedding", "属性から可読な要約文を生成し、多言語モデルで\nベクトル化。属性が変われば トリガ#3 が失効させる"),
             ("VECTOR(384)", "pgvector 型。JP / EN / CN が同一の意味空間を共有する"),
             ("IVFFlat / lists=100", "コサイン距離の Top-K を高速化。空表で作ると\n索引が劣化するため、データ投入後に REINDEX する"),
             ("max_wind_kt は導出値", "経路点から トリガ#2 が自動で上卷するため、\nアプリ側の計算と二重管理にならない")]:
    para(tf, "● " + t, size=12, bold=True, color=NAVY, space_after=2)
    para(tf, "   " + d, size=10, color=GREY, space_after=13, line=1.15)

s = slide("3. 主要テーブル ② track_point", "経路点 = 時空間計算の核。PostGIS ジオメトリと観測時刻を持つ")
code(s, 0.75, 1.32, 7.4, 3.98, [
    "CREATE TABLE track_point (",
    "    id           SERIAL      PRIMARY KEY,",
    "    typhoon_id   INTEGER     NOT NULL",
    "                 REFERENCES typhoon(id) ON DELETE CASCADE,",
    "    agency       VARCHAR(16),        -- CMA / JMA / JTWC / IBTrACS",
    "    obs_time     TIMESTAMPTZ NOT NULL,           -- ▼ 時間軸",
    "    geom         GEOMETRY(POINT, 4326) NOT NULL, -- ▼ 空間軸",
    "    lat          DOUBLE PRECISION NOT NULL,",
    "    lon          DOUBLE PRECISION NOT NULL,",
    "    wind_kt      DOUBLE PRECISION,  pressure_hpa DOUBLE PRECISION,",
    "    grade        VARCHAR(32),       rmw_km       DOUBLE PRECISION,",
    "    move_dir     DOUBLE PRECISION,  move_speed   DOUBLE PRECISION",
    ");",
    "CREATE INDEX ix_track_point_typhoon_id ON track_point (typhoon_id);",
    "CREATE INDEX ix_track_point_obs_time   ON track_point (obs_time);",
    "-- 空間的選択のための GiST 索引",
    "CREATE INDEX ix_track_point_geom ON track_point USING gist (geom);",
], size=10.5, title="DDL — track_point  (135,228 行)")
code(s, 0.75, 5.42, 7.4, 1.35, [
    "-- 主要クエリ Q2 : 時空間選択 → 2306 Khanun が 147 点 (約1か月) 同海域に停滞",
    "SELECT t.intl_id, t.name, COUNT(*) AS hits FROM typhoon t",
    "  JOIN track_point tp ON tp.typhoon_id = t.id",
    " WHERE ST_Intersects(tp.geom, ST_MakeEnvelope(120,20,135,30,4326))",
    "   AND tp.obs_time BETWEEN '2023-07-01' AND '2023-09-30' GROUP BY 1,2;",
], size=10, title="主要クエリ Q2 — 時空間選択 (Spatio-temporal Selection)")
tf = tb(s, 8.4, 1.4, 4.2, 5.3)
para(tf, "設計のポイント", size=15, bold=True, color=BLUE, first=True, space_after=12)
for t, d in [("agency 列で多機関共存", "同一台風に CMA / JMA / JTWC の実況経路が並存する。\n実測 JMA 71,163 ／ CMA 60,498 ／ JTWC 3,567 点"),
             ("統計時は主機関 1 本に正規化", "CMA→JMA→JTWC の優先順で 1 本だけ採用し、\n多重計上を防ぐ"),
             ("geom と lat/lon の二重保持", "geom は空間演算用、lat/lon は API 出力用。\n両者の整合性は トリガ#1 が保証する"),
             ("GiST + B-tree の併用", "ST_Intersects は GiST、obs_time 範囲は B-tree が\n効くため、時空間選択が両軸とも高速"),
             ("ON DELETE CASCADE", "台風を消せば経路点も消える。全子テーブルで統一")]:
    para(tf, "● " + t, size=12, bold=True, color=NAVY, space_after=2)
    para(tf, "   " + d, size=10, color=GREY, space_after=13, line=1.15)

s = slide("3. 主要テーブル ③ secondary_disaster", "受灾情報 = 実際に発生した被害。空間列と意味列を同時に持つ情報層")
code(s, 0.75, 1.32, 7.4, 4.35, [
    "CREATE TABLE secondary_disaster (",
    "    id                SERIAL      PRIMARY KEY,",
    "    typhoon_id        INTEGER     NOT NULL",
    "                      REFERENCES typhoon(id) ON DELETE CASCADE,",
    "    disaster_type     VARCHAR(32) NOT NULL,  -- flood / landslide /",
    "                                           -- storm_surge / casualty",
    "    geom              GEOMETRY(POINT, 4326),  -- ▼ 空間軸",
    "    lat DOUBLE PRECISION,  lon DOUBLE PRECISION,",
    "    event_time        TIMESTAMPTZ,            -- ▼ 時間軸",
    "    casualties        INTEGER,                -- 死傷者数",
    "    economic_loss_usd DOUBLE PRECISION,       -- 経済損失",
    "    description       TEXT,",
    "    region_name       VARCHAR(128),  -- 発信元が報じた地域",
    "    source            VARCHAR(32),   -- GDACS / MEM / FDMA",
    "    source_url        TEXT,",
    "    embedding         VECTOR(384)             -- ▼ 意味軸",
    ");",
    "CREATE INDEX ix_disaster_geom ON secondary_disaster USING gist (geom);",
    "CREATE INDEX ix_disaster_embedding ON secondary_disaster",
    "    USING ivfflat (embedding vector_cosine_ops) WITH (lists = 100);",
], size=10.5, title="DDL — secondary_disaster  (929 行)")
tf = tb(s, 8.4, 1.4, 4.2, 3.0)
para(tf, "情報の二層分離", size=15, bold=True, color=BLUE, first=True, space_after=10)
para(tf, "台風に付随する非構造テキストを、意味的役割で 2 つの対等な表に分ける。",
     size=11, color=GREY, space_after=12, line=1.15)
for t, d, c in [("secondary_disaster ＝ 受灾情報", "既に発生した損失\n死傷・経済損失・洪水 / 土砂 / 高潮\nGDACS 事件・応急管理部・消防庁", ORANGE),
                ("public_info ＝ 公共情報", "当局が公開発表した情報\n警報・避難指示・報道\n気象庁・中央气象台・香港天文台", ACCENT)]:
    para(tf, "■ " + t, size=12, bold=True, color=c, space_after=3)
    para(tf, "   " + d, size=10, color=GREY, space_after=13, line=1.2)
note(s, 8.4, 4.5, 4.2, 1.17,
     "両表とも geom + embedding を持つため、被害情報も警報情報も「時空間でも意味でも」検索できる")
note(s, 0.75, 5.82, 11.8, 0.92,
     "情報を台風へ紐付ける三段階マッチング (load.py)   ① intl_id 一致  →  ② 台風名一致  →  ③ 時空間近傍\n"
     "同一機関が両層を供給することもある。例: GDACS の「事件」は受灾情報、「報道ページ」は公共情報として扱う。")

# ===========================================================================
# 9-11. トリガ
# ===========================================================================
s = slide("4. トリガ ① 空間整合性 — fn_sync_geom()", "緯度経度だけを書いても PostGIS ジオメトリが必ず埋まる")
code(s, 0.75, 1.32, 7.55, 4.85, [
    "CREATE OR REPLACE FUNCTION fn_sync_geom() RETURNS trigger AS $$",
    "BEGIN",
    "    -- lon/lat が揃い、geom が未設定 or 座標がずれている時だけ作り直す",
    "    IF NEW.lon IS NOT NULL AND NEW.lat IS NOT NULL THEN",
    "        IF NEW.geom IS NULL",
    "           OR ST_X(NEW.geom) IS DISTINCT FROM NEW.lon",
    "           OR ST_Y(NEW.geom) IS DISTINCT FROM NEW.lat THEN",
    "            NEW.geom := ST_SetSRID(",
    "                            ST_MakePoint(NEW.lon, NEW.lat), 4326);",
    "        END IF;",
    "    END IF;",
    "    RETURN NEW;",
    "END;",
    "$$ LANGUAGE plpgsql;",
    "",
    "-- 同一関数を 4 つの点テーブルで再利用する (多態トリガ)",
    "CREATE TRIGGER trg_track_point_geom",
    "    BEFORE INSERT OR UPDATE OF lon, lat ON track_point",
    "    FOR EACH ROW EXECUTE FUNCTION fn_sync_geom();",
    "",
    "CREATE TRIGGER trg_secondary_disaster_geom   -- 同様に public_info,",
    "    BEFORE INSERT OR UPDATE OF lon, lat ON secondary_disaster",
    "    FOR EACH ROW EXECUTE FUNCTION fn_sync_geom();   -- landfall にも",
], size=10.5, title="トリガ関数 + トリガ定義")
tf = tb(s, 8.55, 1.4, 4.1, 3.0)
para(tf, "なぜ必要か", size=15, bold=True, color=BLUE, first=True, space_after=10)
para(tf, "geom はすべての空間演算 (ST_Intersects / ST_Contains / GiST 索引) の入口。"
         "ここが 1 行でも欠けると、その台風は空間検索から静かに消える。", size=11, color=GREY,
     space_after=10, line=1.2)
para(tf, "アプリの書き込み漏れ・Navicat からの手入力・外部 SQL 投入 — どの経路から書かれても、"
         "DB が最後の砦として不変条件を守る。", size=11, color=GREY, space_after=0, line=1.2)
code(s, 8.55, 4.42, 4.1, 1.98, [
    "-- 動作確認 (実行済)",
    "INSERT INTO track_point",
    "  (typhoon_id, obs_time,",
    "   lat, lon, wind_kt)",
    "VALUES (84, '2099-01-01',",
    "        25.5, 128.25, 999);",
    "",
    "=> geom = POINT(128.25 25.5)",
], size=10, title="検証")

s = slide("4. トリガ ② 集約ロールアップ — fn_rollup_typhoon_intensity()", "経路点を入れるだけで親台風のサマリ属性が追随する")
code(s, 0.75, 1.32, 7.55, 5.2, [
    "CREATE OR REPLACE FUNCTION fn_rollup_typhoon_intensity()",
    "RETURNS trigger AS $$",
    "BEGIN",
    "    UPDATE typhoon t",
    "       SET max_wind_kt      = GREATEST(t.max_wind_kt,      s.max_wind),",
    "           min_pressure_hpa = LEAST   (t.min_pressure_hpa, s.min_pres),",
    "           start_time       = LEAST   (t.start_time,       s.min_time),",
    "           end_time         = GREATEST(t.end_time,         s.max_time)",
    "      FROM (SELECT typhoon_id,",
    "                   MAX(wind_kt)      AS max_wind,",
    "                   MIN(pressure_hpa) AS min_pres,",
    "                   MIN(obs_time)     AS min_time,",
    "                   MAX(obs_time)     AS max_time",
    "              FROM new_points            -- ▼ 遷移テーブル",
    "             GROUP BY typhoon_id) s",
    "     WHERE t.id = s.typhoon_id;",
    "    RETURN NULL;",
    "END;",
    "$$ LANGUAGE plpgsql;",
    "",
    "CREATE TRIGGER trg_track_point_rollup",
    "    AFTER INSERT ON track_point",
    "    REFERENCING NEW TABLE AS new_points   -- 遷移テーブルを受け取る",
    "    FOR EACH STATEMENT                    -- 行単位ではなく文単位",
    "    EXECUTE FUNCTION fn_rollup_typhoon_intensity();",
], size=10.5, title="トリガ関数 + トリガ定義")
tf = tb(s, 8.55, 1.4, 4.1, 3.0)
para(tf, "設計の要点", size=15, bold=True, color=ACCENT, first=True, space_after=10)
for t, d in [("文レベル + 遷移テーブル", "13 万点の一括投入でも UPDATE は台風数分の\n1 回で済む。行レベルなら 13 万回 UPDATE が\n走り、実用にならない"),
             ("GREATEST / LEAST", "PostgreSQL では NULL を無視するため、\n初回投入 (既存値が NULL) でも正しく動く"),
             ("導出値の一元化", "強度サマリの計算がアプリと DB に\n二重実装されない")]:
    para(tf, "● " + t, size=12, bold=True, color=NAVY, space_after=2)
    para(tf, "   " + d, size=10, color=GREY, space_after=11, line=1.15)
code(s, 8.55, 4.55, 4.1, 1.98, [
    "-- 動作確認 (実行済)",
    "-- 投入前 max_wind_kt=38.9",
    "-- wind 999 / pres 700 の",
    "-- 点を 1 つ INSERT",
    "",
    "=> max_wind_kt      = 999",
    "=> min_pressure_hpa = 700",
    "=> end_time  = 2099-01-01",
], size=10, title="検証")

s = slide("4. トリガ ③ 意味ベクトル失効 — fn_mark_embedding_stale()", "属性が変われば意味ベクトルも自動的に貼り直される")
code(s, 0.75, 1.32, 7.55, 5.35, [
    "CREATE OR REPLACE FUNCTION fn_mark_embedding_stale()",
    "RETURNS trigger AS $$",
    "BEGIN",
    "    NEW.summary_text := NULL;",
    "    NEW.embedding    := NULL;   -- 再埋め込みキューに戻す印",
    "    RETURN NEW;",
    "END;",
    "$$ LANGUAGE plpgsql;",
    "",
    "CREATE TRIGGER trg_typhoon_embedding_stale",
    "    BEFORE UPDATE OF name, name_jp, name_cn, category, season_year,",
    "                     max_wind_kt, min_pressure_hpa,",
    "                     start_time, end_time ON typhoon",
    "    FOR EACH ROW",
    "    WHEN (   OLD.name        IS DISTINCT FROM NEW.name",
    "          OR OLD.category    IS DISTINCT FROM NEW.category",
    "          OR OLD.max_wind_kt IS DISTINCT FROM NEW.max_wind_kt",
    "          OR OLD.min_pressure_hpa IS DISTINCT FROM NEW.min_pressure_hpa",
    "          OR OLD.start_time  IS DISTINCT FROM NEW.start_time",
    "          OR OLD.end_time    IS DISTINCT FROM NEW.end_time)",
    "    EXECUTE FUNCTION fn_mark_embedding_stale();",
    "",
    "-- 受灾情報側も同一関数を再利用 (description / 種別 / 地域 / 死傷者)",
    "CREATE TRIGGER trg_disaster_embedding_stale",
    "    BEFORE UPDATE OF description, disaster_type, region_name, casualties",
    "    ON secondary_disaster FOR EACH ROW WHEN (...)",
    "    EXECUTE FUNCTION fn_mark_embedding_stale();",
], size=10.5, title="トリガ関数 + トリガ定義")
tf = tb(s, 8.55, 1.4, 4.1, 2.6)
para(tf, "トリガ連鎖 (Chain)", size=15, bold=True, color=ORANGE, first=True, space_after=10)
for step, d in [("① 経路点を INSERT", "新しい観測が入る"),
                ("② トリガ② が発火", "typhoon.max_wind_kt を更新"),
                ("③ トリガ③ が発火", "embedding を NULL に失効させる"),
                ("④ embed.py が拾う", "WHERE embedding IS NULL を\n対象に再ベクトル化する")]:
    para(tf, step, size=11.5, bold=True, color=NAVY, space_after=1)
    para(tf, "   " + d, size=10, color=GREY, space_after=9, line=1.15)
note(s, 8.55, 4.2, 4.1, 2.47,
     "WHEN 句が肝。値が実際に変化した UPDATE だけを通すので、冪等な再クロール "
     "(同じ値の書き戻し) では既存のベクトルが無駄に捨てられない。\n\n"
     "検証: 一連の動作確認の後、embedding が NULL の台風は 0 件 — 無駄な失効は発生していない。",
     color=ORANGE)

# ===========================================================================
# 12-14. 関数 + 主要クエリ
# ===========================================================================
s = slide("5. 関数 ① 意味的選択 — fn_semantic_typhoons()", "主要クエリ Q3 : 自然言語 → 384 次元ベクトル → コサイン距離 Top-K")
code(s, 0.75, 1.32, 7.55, 3.6, [
    "CREATE OR REPLACE FUNCTION fn_semantic_typhoons(",
    "    p_qvec vector(384),",
    "    p_k    int DEFAULT 10)",
    "RETURNS TABLE (",
    "    id int, intl_id varchar, name varchar, season_year int,",
    "    max_wind_kt double precision, summary_text text,",
    "    distance double precision",
    ") AS $$",
    "    SELECT t.id, t.intl_id, t.name, t.season_year, t.max_wind_kt,",
    "           t.summary_text,",
    "           (t.embedding <=> p_qvec)::double precision AS distance",
    "      FROM typhoon t",
    "     WHERE t.embedding IS NOT NULL",
    "     ORDER BY t.embedding <=> p_qvec   -- <=> = コサイン距離",
    "     LIMIT p_k;",
    "$$ LANGUAGE sql STABLE",
    "   SET ivfflat.probes = 10;  -- 関数単位の GUC: 再現率を確保する",
], size=10.5, title="関数定義")
code(s, 0.75, 5.05, 7.55, 1.65, [
    "-- 主要クエリ Q3 : 意味的選択 (Semantic Selection)",
    "--   qvec = embed('造成严重洪水和滑坡的强台风')  ← FastAPI 側でベクトル化",
    "SELECT * FROM fn_semantic_typhoons(:qvec, 10);",
    "-- 受灾情報・公共情報にも同形の問合せを並行実行し、3 層横断の意味検索を 1 度に返す",
], size=10, title="主要クエリ Q3 — 呼び出し")
tf = tb(s, 8.55, 1.4, 4.1, 2.4)
para(tf, "設計の要点", size=15, bold=True, color=ORANGE, first=True, space_after=10)
for t, d in [("<=> 演算子", "pgvector のコサイン距離。ORDER BY に置くと\nIVFFlat 索引が使われる"),
             ("関数単位の SET", "STABLE 関数の内部で SET LOCAL は使えない。\n関数属性として GUC を宣言するのが正解"),
             ("probes = 10", "IVFFlat は近似索引。探索クラスタ数を増やして\n再現率を確保する")]:
    para(tf, "● " + t, size=12, bold=True, color=NAVY, space_after=2)
    para(tf, "   " + d, size=10, color=GREY, space_after=11, line=1.15)
code(s, 8.55, 3.95, 4.1, 2.75, [
    "-- 実行結果 (実測)",
    "-- 台風 Gloria(1999) の",
    "-- ベクトルで問合せ",
    "",
    " intl_id  name    year  dist",
    " 9922   Gloria   1999  0.0000",
    " 9607   Gloria   1996  0.0130",
    " 6534   Gloria   1965  0.0168",
    " 7815   Gloria   1978  0.0175",
    " 5226   Gloria   1952  0.0325",
    "",
    "=> 語句一致ではなく意味で",
    "   近い台風が引ける",
], size=9.5, title="検証")

s = slide("5. 関数 ② 時空間的選択 — fn_typhoons_in_bbox()", "主要クエリ Q2 の関数化 : 矩形 × 時間窓を通過した台風を返す")
code(s, 0.75, 1.32, 7.55, 5.4, [
    "CREATE OR REPLACE FUNCTION fn_typhoons_in_bbox(",
    "    p_min_lon double precision, p_min_lat double precision,",
    "    p_max_lon double precision, p_max_lat double precision,",
    "    p_from timestamptz DEFAULT NULL,",
    "    p_to   timestamptz DEFAULT NULL)",
    "RETURNS TABLE (",
    "    id int, intl_id varchar, name varchar, season_year int,",
    "    max_wind_kt double precision,",
    "    hit_points bigint, first_hit timestamptz",
    ") AS $$",
    "    SELECT t.id, t.intl_id, t.name, t.season_year, t.max_wind_kt,",
    "           COUNT(*)         AS hit_points,   -- 矩形内の通過点数",
    "           MIN(tp.obs_time) AS first_hit     -- 初回通過時刻",
    "      FROM typhoon t",
    "      JOIN track_point tp ON tp.typhoon_id = t.id",
    "     WHERE ST_Intersects(tp.geom,            -- GiST 索引が効く",
    "               ST_MakeEnvelope(p_min_lon, p_min_lat,",
    "                               p_max_lon, p_max_lat, 4326))",
    "       AND (p_from IS NULL OR tp.obs_time >= p_from)",
    "       AND (p_to   IS NULL OR tp.obs_time <= p_to)",
    "     GROUP BY t.id, t.intl_id, t.name, t.season_year, t.max_wind_kt",
    "     ORDER BY MIN(tp.obs_time) DESC;",
    "$$ LANGUAGE sql STABLE;",
], size=10.5, title="関数定義")
code(s, 8.55, 1.32, 4.1, 1.75, [
    "-- 呼び出し例",
    "SELECT * FROM",
    " fn_typhoons_in_bbox(",
    "   120, 20, 135, 30,",
    "   '2023-07-01','2023-09-30');",
], size=10, title="呼び出し")
code(s, 8.55, 3.25, 4.1, 1.85, [
    "-- 実行結果 (実測)",
    " intl_id  name        hits",
    " 2312  Kirogi          2",
    " 2313  Yun-Yeung      28",
    " 2311  Haikui         61",
    " 2309  Saola          18",
    " 2306  Khanun        147",
], size=9.5, title="検証")
note(s, 8.55, 5.28, 4.1, 1.42,
     "hit_points が矩形内の滞留の長さを表す。Khanun は 147 点 (約 1 か月) 同海域に停滞しており、"
     "迷走台風であったことが数値から読み取れる。", color=BLUE)

s = slide("5. 関数 ③ 地理影響集計 — fn_region_landfall_count()", "主要クエリ Q5 : 経路を実在の行政境界に帰属させ、上陸頻度として集計する")
code(s, 0.75, 1.32, 7.55, 5.4, [
    "CREATE OR REPLACE FUNCTION fn_region_landfall_count(",
    "    p_region_id int, p_min_year int DEFAULT NULL,",
    "    p_max_year  int DEFAULT NULL) RETURNS bigint AS $$",
    "DECLARE v_level int; v_name varchar; v_count bigint;",
    "BEGIN",
    "    SELECT admin_level, name INTO v_level, v_name",
    "      FROM admin_region WHERE id = p_region_id;",
    "    IF NOT FOUND THEN",
    "        RAISE EXCEPTION '行政区 id=% が存在しません', p_region_id;",
    "    END IF;",
    "",
    "    IF v_level = 0 THEN            -- 国: 非正規化列で高速に数える",
    "        SELECT COUNT(*) INTO v_count FROM landfall l",
    "          JOIN typhoon t ON t.id = l.typhoon_id",
    "         WHERE l.country = v_name",
    "           AND (p_min_year IS NULL OR t.season_year >= p_min_year)",
    "           AND (p_max_year IS NULL OR t.season_year <= p_max_year);",
    "    ELSE                           -- 省 / 地級市: 空間包含で数える",
    "        SELECT COUNT(*) INTO v_count FROM landfall l",
    "          JOIN typhoon t      ON t.id = l.typhoon_id",
    "          JOIN admin_region r ON r.id = p_region_id",
    "         WHERE ST_Contains(r.geom, l.geom)",
    "           AND (p_min_year IS NULL OR t.season_year >= p_min_year)",
    "           AND (p_max_year IS NULL OR t.season_year <= p_max_year);",
    "    END IF;",
    "    RETURN COALESCE(v_count, 0);",
    "END; $$ LANGUAGE plpgsql STABLE;",
], size=10.5, title="関数定義")
tf = tb(s, 8.55, 1.4, 4.1, 2.0)
para(tf, "なぜ粒度で分岐するか", size=15, bold=True, color=ACCENT, first=True, space_after=9)
para(tf, "上陸イベントは「最も細かい行政区」1 つに紐付けて保存される。国別集計では非正規化した "
         "country 列を使う方が圧倒的に速い。", size=10.5, color=GREY, space_after=8, line=1.2)
para(tf, "一方、省・地級市の集計を紐付け先に頼ると粒度違いの上陸点を取りこぼすため、"
         "ST_Contains による空間包含で数え直す。", size=10.5, color=GREY, space_after=0, line=1.2)
code(s, 8.55, 3.6, 4.1, 3.1, [
    "-- 実行結果 (実測)",
    "SELECT fn_region_",
    "  landfall_count(3);",
    "=> 600   -- China (国)",
    "",
    "SELECT fn_region_",
    "  landfall_count(173);",
    "=> 217   -- 広東省",
    "",
    "-- 省別 (実測)",
    " 広東 217 / 海南 121",
    " 福建 111 / 浙江  44",
    " 広西  41 / 山東  16",
], size=9.5, title="検証")

# ===========================================================================
# 15. 主要クエリ Q4 (結合) + まとめ
# ===========================================================================
s = slide("6. 主要クエリ Q4 — 時空間 × 意味 の結合 ／ まとめ", "先に時空間で候補を絞り、その中を意味距離で並べ替える — 本システムの中心的な問合せ")
code(s, 0.75, 1.32, 6.9, 4.05, [
    "-- 「高潮による沿岸被害」に意味が近い台風のうち、",
    "-- 2023 年に東シナ海を通過したものを上位から",
    "SELECT t.intl_id, t.name, t.season_year,",
    "       t.embedding <=> :qvec AS distance",
    "  FROM typhoon t",
    " WHERE t.id IN (              -- ① 時空間選択で候補を作る",
    "        SELECT tp.typhoon_id FROM track_point tp",
    "         WHERE ST_Intersects(tp.geom,",
    "                 ST_MakeEnvelope(120, 20, 135, 30, 4326))",
    "           AND tp.obs_time BETWEEN '2023-01-01'",
    "                               AND '2023-12-31')",
    "   AND t.embedding IS NOT NULL",
    " ORDER BY t.embedding <=> :qvec     -- ② 意味距離で順序付け",
    " LIMIT 10;",
    "",
    "-- ① の時空間フィルタで候補は 9 件に絞られた (実測)",
    "-- ② その 9 件だけに重い意味計算を適用する",
], size=10.5, title="SQL — /search/hybrid")
note(s, 0.75, 5.5, 6.9, 1.22,
     "順序が重要。時空間 (GiST) → 意味 (IVFFlat) の順なら、重い意味計算を絞り込み後の少数行にだけ適用できる。"
     "逆順では Top-K が時空間条件で全滅し、結果が空になりうる。",
     color=RGBColor(0x8E, 0x44, 0xAD))

rect(s, 7.95, 1.32, 4.6, 0.42, fill=NAVY)
htf = tb(s, 8.15, 1.42, 4.2, 0.3)
para(htf, "検証結果 (MMKB 本番データ / 2026-07-19)", size=11.5, bold=True, color=WHITE, first=True, space_after=0)
rows = [
    ("トリガ①", "lat/lon のみ INSERT → geom 自動生成"),
    ("トリガ②", "経路点 1 点追加 → 親台風 38.9 → 999 に追随"),
    ("トリガ③", "属性更新 → embedding が NULL に失効"),
    ("関数①", "Top-5 を距離順に返却 (0.0000 / 0.0130 …)"),
    ("関数②", "2023 夏の通過台風 5 件を通過点数つきで返却"),
    ("関数③", "China=600 ／ 広東省=217 ／ 同 2000-20=63"),
    ("クエリ", "Q1〜Q5 すべて本番データで期待どおり"),
]
y = 1.74
for i, (a, b) in enumerate(rows):
    bg = WHITE if i % 2 == 0 else RGBColor(0xF6, 0xF8, 0xFB)
    rect(s, 7.95, y, 4.6, 0.44, fill=bg, line=BORDER)
    t1 = tb(s, 8.1, y + 0.12, 0.95, 0.28)
    para(t1, a, size=10, bold=True, color=NAVY, first=True, space_after=0)
    t2 = tb(s, 9.05, y + 0.13, 3.0, 0.28)
    para(t2, b, size=8.5, color=GREY, first=True, space_after=0)
    t3 = tb(s, 12.1, y + 0.12, 0.4, 0.28)
    para(t3, "✔", size=11, bold=True, color=RGBColor(0x1E, 0x8E, 0x4A), first=True, space_after=0)
    y += 0.44

note(s, 7.95, 4.9, 4.6, 0.95,
     "トリガ検証はトランザクション内で実行し ROLLBACK 済み。本番データは変更していない。"
     "全定義は docs/db_objects.sql に冪等スクリプトとして保存。")
rect(s, 7.95, 5.98, 4.6, 0.74, fill=NAVY)
tf = tb(s, 8.15, 6.14, 4.2, 0.5)
para(tf, "「時空間で絞り、意味で結ぶ」", size=13.5, bold=True, color=WHITE, first=True,
     align=PP_ALIGN.CENTER, space_after=2)
para(tf, "分散マルチメディア情報を知識として扱う一実装", size=9.5,
     color=RGBColor(0x9F, 0xC4, 0xE2), align=PP_ALIGN.CENTER, space_after=0)

prs.save(OUT)
io.open(r"C:\Users\49879\AppData\Local\Temp\claude\d--Project-Musashino-ADVANCED-MMKB-2026\a085c82e-07ab-4fd5-9340-d8318cc6bb9d\scratchpad\ppt.txt",
        "w", encoding="utf-8").write(
    "saved %s\nslides=%d\n\nWARNINGS:\n%s" % (
        OUT, len(prs.slides._sldIdLst), "\n".join(WARN) or "(none)"))
